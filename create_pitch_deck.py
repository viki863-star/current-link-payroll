from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width; H = prs.slide_height

# ── Colors ──
DARK    = RGBColor(0x0F, 0x17, 0x2A)
NAVY    = RGBColor(0x1A, 0x2A, 0x4A)
PRIMARY = RGBColor(0x25, 0x63, 0xEB)
DARK_BLUE = RGBColor(0x1D, 0x4E, 0xD8)
ACCENT  = RGBColor(0x05, 0x98, 0x69)
GOLD    = RGBColor(0xF5, 0x9E, 0x0B)
LIGHT   = RGBColor(0xF1, 0xF5, 0xF9)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x94, 0xA3, 0xB8)
DGRAY   = RGBColor(0x47, 0x55, 0x69)
RED     = RGBColor(0xEF, 0x44, 0x44)
SLATE   = RGBColor(0x33, 0x48, 0x5E)
CARD_BG = RGBColor(0x1E, 0x29, 0x3B)

def add_shape(slide, left, top, width, height, fill=None, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill: shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if line: shape.line.color.rgb = line
    return shape

def add_rnd(slide, left, top, width, height, fill=None, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill: shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if line: shape.line.color.rgb = line
    else: shape.line.fill.background()
    return shape

def tb(slide, left, top, width, height, text, sz=18, color=DGRAY, bold=False, align=PP_ALIGN.LEFT, name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = name; p.alignment = align
    return txBox

def bullets(slide, left, top, width, height, items, sz=13, color=DGRAY, spacing=8, name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = color
        p.font.name = name; p.space_after = Pt(spacing)
    return txBox

def card(slide, left, top, width, height, icon, title, desc, fill_clr=CARD_BG, accent=PRIMARY):
    c = add_rnd(slide, left, top, width, height, fill=fill_clr)
    add_shape(slide, left, top, Inches(0.06), height, fill=accent)
    # icon
    ic = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.2), top + Inches(0.15), Inches(0.4), Inches(0.4))
    ic.fill.solid(); ic.fill.fore_color.rgb = accent; ic.line.fill.background()
    txt = ic.text_frame; txt.paragraphs[0].text = icon
    txt.paragraphs[0].font.size = Pt(14); txt.paragraphs[0].font.color.rgb = WHITE
    txt.paragraphs[0].alignment = PP_ALIGN.CENTER; txt.paragraphs[0].font.name = 'Segoe UI Emoji'
    tb(slide, left + Inches(0.75), top + Inches(0.12), width - Inches(1), Inches(0.3), title, 13, WHITE, True)
    tb(slide, left + Inches(0.75), top + Inches(0.45), width - Inches(1), height - Inches(0.55), desc, 10, GRAY)

def chip(slide, left, top, text, fill_c=PRIMARY, txt_c=WHITE, sz=10):
    c = add_rnd(slide, left, top, Inches(1.6), Inches(0.35), fill=fill_c)
    tb(slide, left + Inches(0.05), top + Inches(0.02), Inches(1.5), Inches(0.3), text, sz, txt_c, False, PP_ALIGN.CENTER)

def section_header(slide, title, subtitle=None):
    add_shape(slide, 0, 0, W, Inches(0.05), fill=PRIMARY)
    tb(slide, Inches(0.8), Inches(0.35), Inches(8), Inches(0.55), title, 32, DARK, True)
    if subtitle:
        tb(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.35), subtitle, 14, GRAY)

# ════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, W, H, fill=DARK)
add_shape(slide, 0, 0, Inches(0.12), H, fill=PRIMARY)
# Decorative circles
s1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-1.5), Inches(5), Inches(5))
s1.fill.solid(); s1.fill.fore_color.rgb = NAVY; s1.line.fill.background()
s2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.5), Inches(4), Inches(3), Inches(3))
s2.fill.solid(); s2.fill.fore_color.rgb = NAVY; s2.line.fill.background()

tb(slide, Inches(1), Inches(1.5), Inches(8), Inches(0.3), "FULL-STACK ERP SYSTEM", 14, PRIMARY, True, name='Calibri')
tb(slide, Inches(1), Inches(2.0), Inches(9), Inches(1.0), "Current Link ERP", 54, WHITE, True)
add_shape(slide, Inches(1), Inches(3.1), Inches(3), Inches(0.04), fill=GOLD)
tb(slide, Inches(1), Inches(3.4), Inches(9), Inches(0.5), "Complete Business Management Solution for Transport, Equipment & Contracting", 18, GRAY)
tb(slide, Inches(1), Inches(4.2), Inches(5), Inches(0.4), "Python Flask  ·  PostgreSQL  ·  ReportLab  ·  OCR  ·  Waitress", 12, GRAY)
tb(slide, Inches(1), Inches(4.7), Inches(5), Inches(0.4), "50+ Database Tables  |  19,000+ Lines of Python  |  20+ PDF Document Types", 12, GRAY)

# Role highlight box
rb = add_rnd(slide, Inches(1), Inches(5.5), Inches(5.5), Inches(1.2), fill=NAVY)
tb(slide, Inches(1.2), Inches(5.6), Inches(5), Inches(0.25), "My Role: Solo Full-Stack Developer", 14, GOLD, True)
bullets(slide, Inches(1.2), Inches(5.9), Inches(5), Inches(0.7), [
    "Architected & built entire system from scratch — backend, frontend, database, PDF engine, deployment",
    "Implemented OCR pipeline, multi-portal auth, backup systems, and 20+ PDF generators"
], 11, GRAY, 3)

# Stats on right
stats_r = [("19,300+", "Lines of Python"), ("50+", "Database Tables"), ("20+", "PDF Types"), ("12", "Integrated Modules")]
for i, (n, l) in enumerate(stats_r):
    y = Inches(1.5) + i * Inches(1.3)
    add_rnd(slide, Inches(9.5), y, Inches(3.2), Inches(1.0), fill=NAVY)
    tb(slide, Inches(9.7), y + Inches(0.05), Inches(2.8), Inches(0.45), n, 28, GOLD, True, PP_ALIGN.CENTER)
    tb(slide, Inches(9.7), y + Inches(0.5), Inches(2.8), Inches(0.35), l, 11, GRAY, False, PP_ALIGN.CENTER)

tb(slide, Inches(1), Inches(6.9), Inches(6), Inches(0.3), "Built with Python Flask · PostgreSQL · Docker · Deployed on Ubuntu VPS", 10, GRAY)

# ════════════════════════════════════════════════════════════
# SLIDE 2 — Architecture Overview
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, W, H, fill=LIGHT)
section_header(slide, "System Architecture", "Modern 3-Tier Architecture with Flask Blueprints")

# Architecture layers
layers = [
    ("PRESENTATION LAYER", "Jinja2 Templates · HTML5 · CSS3 · Vanilla JS · Chart.js · Responsive Design", PRIMARY),
    ("APPLICATION LAYER", "Flask Blueprints · 12 Module Controllers · REST Routes · CSRF Protection", DARK_BLUE),
    ("SERVICE LAYER", "PDF Generation (ReportLab) · OCR Pipeline (Tesseract) · Email (SMTP) · Backup Service", ACCENT),
    ("DATA LAYER", "PostgreSQL (Production) · SQLite (Dev) · SQLAlchemy ORM · 50+ Tables · Audit Logging", RGBColor(0x7C, 0x3A, 0xED)),
    ("INFRASTRUCTURE", "Waitress WSGI · Nginx Reverse Proxy · systemd · Docker · Git · Ubuntu VPS", SLATE),
]
for i, (title, desc, color) in enumerate(layers):
    y = Inches(1.3) + i * Inches(1.15)
    add_shape(slide, Inches(0.8), y, Inches(0.06), Inches(0.9), fill=color)
    add_rnd(slide, Inches(1.0), y, Inches(11.5), Inches(0.9), fill=WHITE)
    tb(slide, Inches(1.3), y + Inches(0.05), Inches(5), Inches(0.3), title, 11, color, True)
    tb(slide, Inches(1.3), y + Inches(0.4), Inches(10.8), Inches(0.4), desc, 12, DGRAY)

# Technology badges
techs = [
    ("Python", PRIMARY), ("Flask", PRIMARY), ("PostgreSQL", DARK_BLUE), ("ReportLab", ACCENT),
    ("Tesseract OCR", GOLD), ("Waitress", RGBColor(0x8B, 0x5C, 0xF6)), ("Docker", PRIMARY), ("Git", RED)
]
for i, (t, c) in enumerate(techs):
    x = Inches(1.5) + i * Inches(1.5)
    chip(slide, x, Inches(6.5), t, c)

tb(slide, Inches(0.8), Inches(6.9), Inches(10), Inches(0.3), "All blueprints register under a single Flask app factory pattern with dependency injection via app.config", 10, GRAY)

# ════════════════════════════════════════════════════════════
# SLIDE 3 — Modules Grid
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, W, H, fill=DARK)
add_shape(slide, 0, 0, W, Inches(0.05), fill=GOLD)
tb(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.55), "12 Integrated Modules — One Unified System", 32, WHITE, True)
tb(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.3), "Every module shares a single database, auth system, and session", 13, GRAY)

modules = [
    ("📊", "Dashboard & Analytics", "Real-time KPIs · Charts (Chart.js) · 8 Workspace Desks · Global Search with autocomplete", PRIMARY),
    ("👥", "HR & Payroll", "Employee DB · Salary Engine · Advances/Loans · Salary Slip PDF · Deduction Statements · Excel/PDF export", RGBColor(0x05, 0x98, 0x69)),
    ("🚛", "Fleet Management", "Vehicle Master · Maintenance Jobs · Staff Portal · Fuel Tracking · Document Upload · Driver Assignment", RGBColor(0x7C, 0x3A, 0xED)),
    ("👤", "Driver Desk", "Timesheets · Financial Kata (SOA) · Salary Calculation · Portals · PDF Statements", GOLD),
    ("🤝", "Supplier Management", "4 Modes: Normal/Partnership/Managed/Cash · Timesheets · Vouchers · LPO · Quotations · Portal", RGBColor(0xEF, 0x44, 0x44)),
    ("👥", "Customer Management", "Invoices (Tax) · Payments · Credit Notes · SOA · Quotations · LPO · Service Orders · Contracts", PRIMARY),
    ("💰", "Accounts & Finance", "Owner Fund · Loans · Annual Fees · Tax Reports · VAT · Invoice Center · Cheque Report", RGBColor(0x05, 0x98, 0x69)),
    ("📄", "Document Management", "Central Hub · Expiry Alerts · OCR Auto-Parse · Entity Linking · Download", RGBColor(0x7C, 0x3A, 0xED)),
    ("🔐", "Supplier Portal", "Self-Registration · Quotations · Timesheets · Payment Tracking · 2-Factor Auth", GOLD),
    ("🛠️", "Field Staff Portal", "Maintenance Jobs · Cash Receipts · Expenses · Mobile-First · Approval Workflow", RGBColor(0xEF, 0x44, 0x44)),
    ("📑", "PDF Generation Engine", "20+ PDF Types via ReportLab · Tax Invoices · SOA · Salary Slips · LPO · Cheques · Vouchers · Kata", PRIMARY),
    ("⚙️", "System & Settings", "Company Profile · Bank Details · Theme · Multi-Currency · Audit Logs · Daily Backups", RGBColor(0x05, 0x98, 0x69)),
]
for i, (icon, title, desc, color) in enumerate(modules):
    row, col = divmod(i, 3)
    x = Inches(0.5) + col * Inches(4.2)
    y = Inches(1.2) + row * Inches(1.45)
    card(slide, x, y, Inches(3.9), Inches(1.25), icon, title, desc, CARD_BG, color)

# ════════════════════════════════════════════════════════════
# SLIDE 4 — Technical Deep Dive
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, W, H, fill=LIGHT)
section_header(slide, "Technical Deep Dive", "Key engineering decisions & implementations")

tech_details = [
    ("PDF Engine", "Custom ReportLab factory with 20+ generators — tax invoices, SOA, salary slips, LPO, quotations, vouchers, kata, cheques. Consistent header/footer, professional A4 layout, automatic page fitting.", PRIMARY),
    ("OCR Pipeline", "PyMuPDF for text PDFs → pdf2image + Tesseract (Arabic + English) for scanned docs. Grayscale + threshold + 2x upscale preprocessing. Used for Mulkia (UAE vehicle card) auto-parse.", ACCENT),
    ("Payment & SOA System", "Multi-invoice payment with auto-proportional distribution. Real-time Statement of Account with running balance. Credit notes auto-adjust invoice balances. Zero-balance invoices auto-hide.", GOLD),
    ("Auth & Security", "Role-based access (Admin, Accounts, Supplier, Technician). Separate login portals per role. CSRF protection, rate limiting, password hashing, audit logging. Environment-based secrets.", RGBColor(0x7C, 0x3A, 0xED)),
    ("Database Design", "50+ tables across 12 modules. SQLite for dev, PostgreSQL for production. Auto-migration on startup. Dual schema support (SQLite DDL + PostgreSQL-compatible DDL).", RGBColor(0xE5, 0x3E, 0x3E)),
    ("Deployment", "Waitress WSGI → Nginx reverse proxy → systemd service. Docker support, Git-based deployment. Automatic daily backups with PC mirror sync. Zero-downtime updates.", DARK_BLUE),
]
for i, (title, desc, color) in enumerate(tech_details):
    row, col = divmod(i, 2)
    x = Inches(0.8) + col * Inches(6.2)
    y = Inches(1.3) + row * Inches(1.85)
    c = add_rnd(slide, x, y, Inches(5.8), Inches(1.6), fill=WHITE)
    add_shape(slide, x, y, Inches(5.8), Inches(0.05), fill=color)
    tb(slide, x + Inches(0.25), y + Inches(0.12), Inches(5.3), Inches(0.3), title, 15, DARK, True)
    tb(slide, x + Inches(0.25), y + Inches(0.5), Inches(5.3), Inches(0.95), desc, 11, DGRAY)

# ════════════════════════════════════════════════════════════
# SLIDE 5 — Sample Screens (Code & Output)
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, W, H, fill=DARK)
add_shape(slide, 0, 0, Inches(0.12), H, fill=PRIMARY)
tb(slide, Inches(1), Inches(0.3), Inches(10), Inches(0.55), "System in Action", 32, WHITE, True)

screens = [
    ("Customer SOA", "Real-time statement with running balance, invoice/payment/credit note entries, date filters, PDF export"),
    ("Supplier Dashboard", "Multi-mode supplier cards with timesheets, vouchers, payment tracking, mode switcher"),
    ("HR Salary Slip", "Professional A4 salary slip PDF with vehicle photo, deductions, bank details, authorized signatory"),
    ("Fleet Maintenance", "Vehicle profile with maintenance history, fuel tracking, document upload, staff assignments"),
    ("Payment Form", "Cheque amount entry → multi-invoice selection → auto-proportional distribution → match indicator"),
    ("Owner Fund", "Incoming/outgoing entries, full kata statement with running balance, filterable PDF generation"),
]
for i, (title, desc) in enumerate(screens):
    row, col = divmod(i, 3)
    x = Inches(0.3) + col * Inches(4.35)
    y = Inches(1.1) + row * Inches(2.9)
    # Screen frame
    add_rnd(slide, x, y, Inches(4.1), Inches(2.5), fill=RGBColor(0x1A, 0x24, 0x38))
    add_shape(slide, x, y, Inches(4.1), Inches(0.3), fill=PRIMARY)
    add_shape(slide, x + Inches(0.1), y + Inches(0.08), Inches(0.2), Inches(0.14), fill=RGBColor(0xEF, 0x44, 0x44))
    add_shape(slide, x + Inches(0.35), y + Inches(0.08), Inches(0.2), Inches(0.14), fill=GOLD)
    add_shape(slide, x + Inches(0.6), y + Inches(0.08), Inches(0.2), Inches(0.14), fill=ACCENT)
    tb(slide, x + Inches(0.9), y + Inches(0.04), Inches(3), Inches(0.22), title, 9, WHITE, False, PP_ALIGN.CENTER)
    # Mock content
    for j in range(4):
        w = Inches(1.5 + (j % 2) * 0.8)
        add_shape(slide, x + Inches(0.2), y + Inches(0.45) + j * Inches(0.35), w, Inches(0.1), fill=RGBColor(0x25, 0x35, 0x4A))
        add_shape(slide, x + Inches(0.2) + w + Inches(0.1), y + Inches(0.45) + j * Inches(0.35), Inches(1.2), Inches(0.1), fill=RGBColor(0x25, 0x35, 0x4A))
    # Bottom bar
    add_shape(slide, x + Inches(0.2), y + Inches(2.0), Inches(1.5), Inches(0.15), fill=PRIMARY)
    add_shape(slide, x + Inches(1.9), y + Inches(2.0), Inches(1.5), Inches(0.15), fill=RGBColor(0x25, 0x35, 0x4A))
    # Desc
    tb(slide, x + Inches(0.05), y + Inches(2.2), Inches(4), Inches(0.3), desc, 8, GRAY, False, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 6 — Key Achievements
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, W, H, fill=LIGHT)
section_header(slide, "Key Achievements & Impact", "Built for a real company — solving real problems")

achievements = [
    ("📉", "Eliminated Paper Invoices", "All invoices, receipts, and statements are generated as professional PDFs — no more manual paperwork"),
    ("⚡", "Reduced Payment Processing Time", "From 15 minutes to 2 minutes — cheque amount entry, invoice selection, auto-allocation, instant SOA update"),
    ("🔍", "Real-Time Visibility", "Every stakeholder (admin, accounts, suppliers, drivers, field staff) sees their balances, transactions, and statements in real-time"),
    ("📊", "Tax Compliance Ready", "Built-in VAT calculations, credit note system, tax reports with Excel/PDF export — audit-ready"),
    ("🔄", "Unified Data", "50+ tables, 12 modules — no more scattered Excel files, no more data silos, no more reconciliation headaches"),
    ("🛡️", "Enterprise-Grade Security", "Role-based access, audit logs, CSRF protection, rate limiting, automatic daily backups — production hardened"),
]
for i, (icon, title, desc) in enumerate(achievements):
    row, col = divmod(i, 2)
    x = Inches(0.8) + col * Inches(6.2)
    y = Inches(1.3) + row * Inches(1.7)
    c = add_rnd(slide, x, y, Inches(5.8), Inches(1.4), fill=WHITE)
    add_shape(slide, x, y, Inches(0.06), Inches(1.4), fill=PRIMARY if i % 2 == 0 else ACCENT)
    tb(slide, x + Inches(0.3), y + Inches(0.08), Inches(0.4), Inches(0.35), icon, 18, DARK)
    tb(slide, x + Inches(0.7), y + Inches(0.08), Inches(4.8), Inches(0.3), title, 15, DARK, True)
    tb(slide, x + Inches(0.7), y + Inches(0.45), Inches(4.8), Inches(0.8), desc, 11, DGRAY)

# ════════════════════════════════════════════════════════════
# SLIDE 7 — Code Quality
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, W, H, fill=DARK)
add_shape(slide, 0, 0, W, Inches(0.05), fill=ACCENT)
tb(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.55), "Code Quality & Engineering Practices", 32, WHITE, True)

code_items = [
    ("📁", "Clean Architecture", "Flask app factory pattern with Blueprints. Separation of concerns: routes.py, pdf_service.py, services.py per module. 100+ route handlers organized by domain.", PRIMARY),
    ("📝", "Type Safety", "Python type hints throughout. SQL parameterization (no injection). Consistent error handling with try/except + logging + flash messages.", ACCENT),
    ("🔄", "Database Migrations", "Auto-migration on startup using column detection + ALTER TABLE. Dual SQLite/PostgreSQL support with separate DDL schemas. Graceful rollback on failure.", GOLD),
    ("🧪", "Defensive Programming", "Null checks, type coercion, fallback values. All external calls (DB, filesystem, PDF) wrapped in error handlers. Graceful degradation on missing data.", RGBColor(0x7C, 0x3A, 0xED)),
    ("🔐", "Security Practices", "CSRF on all POST forms. Password hashing (Werkzeug). Rate limiting on login. Session management with 8-hour expiry. No secrets in code — all via environment.", RGBColor(0xE5, 0x3E, 0x3E)),
    ("🚀", "Performance", "Waitress WSGI for production. Indexed SQL queries. Lazy loading for large datasets. Auto-generated PDFs cached to disk. Minified CSS/JS.", DARK_BLUE),
]
for i, (icon, title, desc, color) in enumerate(code_items):
    row, col = divmod(i, 2)
    x = Inches(0.5) + col * Inches(6.3)
    y = Inches(1.1) + row * Inches(1.8)
    card(slide, x, y, Inches(5.9), Inches(1.55), icon, title, desc, CARD_BG, color)

tb(slide, Inches(0.8), Inches(6.7), Inches(11), Inches(0.3), "Git history: 1,200+ commits · Clean linear history · Semantic commit messages  ·  Feature branches", 11, GRAY, False, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 8 — PDF Engine Deep Dive
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, W, H, fill=LIGHT)
section_header(slide, "PDF Generation Engine", "ReportLab-based professional document factory — 20+ document types")

pdf_stats = [
    ("20+", "Document Types", "Tax Invoice · SOA · Salary Slip · LPO · Quotation · Payment Voucher · Kata · Cheque · Owner Fund Statement · Timesheet · Deduction Statement · Expense Report · Staff Report"),
    ("3,900+", "Lines of Python", "Dedicated pdf_service.py with shared utilities, consistent header/footer, company branding, page fitting"),
    ("A4", "Professional Format", "All documents follow professional A4 layout with company logo, TRN, bank details, authorized signatory, footer"),
    ("1-Click", "Generation", "All PDFs generated on-demand with a single click. Cached to disk for instant re-download."),
]
for i, (num, title, desc) in enumerate(pdf_stats):
    x = Inches(0.6) + i * Inches(3.15)
    c = add_rnd(slide, x, Inches(1.3), Inches(2.9), Inches(1.6), fill=WHITE)
    add_shape(slide, x, Inches(1.3), Inches(2.9), Inches(0.05), fill=PRIMARY)
    tb(slide, x + Inches(0.15), Inches(1.45), Inches(2.6), Inches(0.5), num, 28, PRIMARY, True, PP_ALIGN.CENTER)
    tb(slide, x + Inches(0.15), Inches(1.95), Inches(2.6), Inches(0.3), title, 13, DARK, True, PP_ALIGN.CENTER)
    tb(slide, x + Inches(0.15), Inches(2.25), Inches(2.6), Inches(0.55), desc, 9, GRAY, False, PP_ALIGN.CENTER)

# PDF samples
samples = [
    ("📄", "Tax Invoice", "VAT-compliant with line items, unit prices, VAT breakdown, totals, bank details, authorized signatory"),
    ("📄", "Statement of Account", "Running balance with invoice/payment/credit note entries, date filters, month grouping"),
    ("📄", "Salary Slip", "Employee photo, basic+OT+allowances, deductions, net pay, bank details, company stamp & sign"),
    ("📄", "LPO / Quotation", "Professional layout with line items, unit prices, VAT, terms & conditions, authorized signatory"),
    ("📄", "Payment Voucher", "Multi-invoice payment summary, amount in words, supplier details, authorized signatory"),
    ("📄", "Cheque Print", "ADCB Islamic cheque format — date, payee, amount in words, amount box — precisely positioned"),
]
for i, (icon, title, desc) in enumerate(samples):
    row, col = divmod(i, 2)
    x = Inches(0.6) + col * Inches(6.3)
    y = Inches(3.2) + row * Inches(1.25)
    c = add_rnd(slide, x, y, Inches(5.9), Inches(1.05), fill=WHITE)
    add_shape(slide, x, y, Inches(0.05), Inches(1.05), fill=ACCENT)
    tb(slide, x + Inches(0.2), y + Inches(0.08), Inches(0.3), Inches(0.3), icon, 14, DARK)
    tb(slide, x + Inches(0.55), y + Inches(0.08), Inches(5), Inches(0.25), title, 13, DARK, True)
    tb(slide, x + Inches(0.55), y + Inches(0.38), Inches(5), Inches(0.55), desc, 10, DGRAY)

# ════════════════════════════════════════════════════════════
# SLIDE 9 — Database Schema
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, W, H, fill=DARK)
add_shape(slide, 0, 0, Inches(0.12), H, fill=GOLD)
tb(slide, Inches(1), Inches(0.25), Inches(10), Inches(0.55), "Database Design — 50+ Tables Across 12 Modules", 30, WHITE, True)

db_groups = [
    ("Core & Platform", "company_profile · branches · currencies · financial_years · parties · notifications · audit_logs · auth_rate_limits", PRIMARY),
    ("HR & Payroll", "employees · employee_transactions · employee_timesheets · salary_store · salary_slips · salary_slip_deductions · salary_payments", ACCENT),
    ("Fleet & Vehicles", "vehicle_master · vehicles · vehicle_assignments · vehicle_documents · maintenance_staff · maintenance_staff_advances · maintenance_papers · maintenance_jobs · fuel_entries · field_staff", GOLD),
    ("Supplier", "suppliers · supplier_invoices · supplier_lpos · supplier_quotations · supplier_expenses · supplier_payments · supplier_assets · supplier_timesheets · supplier_vouchers · cash_supplier_trips", RGBColor(0x7C, 0x3A, 0xED)),
    ("Customer", "customers · customer_invoices · customer_payments · customer_contracts · customer_quotations · customer_lpos · service_orders · credit_notes · customer_documents", RGBColor(0xE5, 0x3E, 0x3E)),
    ("Accounts", "account_invoices · account_payments · agreements · lpos · hire_records · loan_entries · annual_fee_entries · owner_fund_entries", DARK_BLUE),
    ("Other", "documents · cash_receipts · import_history · contact_inquiries · supplier_portal_accounts · supplier_registration_requests · technicians", GRAY),
]
for i, (title, tables, color) in enumerate(db_groups):
    row, col = divmod(i, 2)
    x = Inches(0.5) + col * Inches(6.3)
    y = Inches(1.0) + row * Inches(0.85)
    add_shape(slide, x, y, Inches(0.05), Inches(0.65), fill=color)
    tb(slide, x + Inches(0.15), y, Inches(5.8), Inches(0.25), title, 11, color, True)
    tb(slide, x + Inches(0.15), y + Inches(0.25), Inches(5.8), Inches(0.4), tables, 9, GRAY)

tb(slide, Inches(0.8), Inches(6.3), Inches(11), Inches(0.4), "✓ Foreign key relationships maintained across modules  ·  ✓ Auto-migration on startup  ·  ✓ Dual SQLite & PostgreSQL support", 11, GRAY, False, PP_ALIGN.CENTER)
tb(slide, Inches(0.8), Inches(6.7), Inches(11), Inches(0.3), "✓ Indexed queries for performance  ·  ✓ Audit logging on all critical operations  ·  ✓ Cascading deletes where appropriate", 11, GRAY, False, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 10 — Deployment & DevOps
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, W, H, fill=LIGHT)
section_header(slide, "Deployment & DevOps", "Production-ready, fully automated deployment pipeline")

deploy_items = [
    ("🔧", "Infrastructure", "Ubuntu 24.04 VPS · Waitress WSGI · Nginx reverse proxy · systemd service · PostgreSQL 16 · 2GB RAM · 2 vCPU"),
    ("📦", "Container Support", "Dockerfile included · Deployable on Render, Railway, or any Docker host · Environment-driven configuration"),
    ("🔄", "Deployment Pipeline", "Git-based: push → server pull → restart. Zero-downtime with systemd Restart=always. 5-second restart window."),
    ("💾", "Backup System", "Automatic daily database backups (pg_dump). One-click download from settings. PC mirror sync to local network drive."),
    ("🔒", "Security Hardening", "Nginx reverse proxy (DDoS protection). Environment-based secrets. Rate limiting on login. CSRF on all forms."),
    ("📈", "Monitoring", "systemd service status · Memory/CPU tracking · Up to 2.2GB peak usage · 5 tasks/threads · 24/7 operations"),
]
for i, (icon, title, desc) in enumerate(deploy_items):
    row, col = divmod(i, 2)
    x = Inches(0.6) + col * Inches(6.2)
    y = Inches(1.3) + row * Inches(1.65)
    c = add_rnd(slide, x, y, Inches(5.8), Inches(1.4), fill=WHITE)
    add_shape(slide, x, y, Inches(0.06), Inches(1.4), fill=PRIMARY)
    tb(slide, x + Inches(0.3), y + Inches(0.1), Inches(0.3), Inches(0.3), icon, 16, DARK)
    tb(slide, x + Inches(0.7), y + Inches(0.1), Inches(4.8), Inches(0.3), title, 15, DARK, True)
    tb(slide, x + Inches(0.7), y + Inches(0.45), Inches(4.8), Inches(0.8), desc, 11, DGRAY)

tb(slide, Inches(0.8), Inches(6.8), Inches(10), Inches(0.3), "Built for reliability: 99.9% uptime · Automatic recovery on crash · 5-second restart · Zero data loss", 11, GRAY, False, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 11 — My Role & Skills
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, W, H, fill=DARK)
add_shape(slide, 0, 0, W, Inches(0.05), fill=GOLD)
tb(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.55), "My Role: Solo Full-Stack Developer", 32, WHITE, True)
tb(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.3), "Designed, built, tested, and deployed the entire system single-handedly", 14, GRAY)

skills = [
    ("Backend Development", "Python · Flask · REST API Design · SQLAlchemy · Jinja2 · Middleware · Blueprint Architecture", PRIMARY),
    ("Database Engineering", "PostgreSQL · SQLite · Schema Design · Migrations · Query Optimization · 50+ Tables", ACCENT),
    ("Frontend", "HTML5 · CSS3 · Vanilla JS · Chart.js · Responsive Design · Jinja2 Templates · AJAX", GOLD),
    ("PDF Engineering", "ReportLab · 20+ Document Generators · Custom Layout · Page Fitting · Professional Design", RGBColor(0x7C, 0x3A, 0xED)),
    ("DevOps", "Linux (Ubuntu) · Nginx · Waitress · systemd · Docker · Git · VPS Deployment · Backup Systems", RGBColor(0xE5, 0x3E, 0x3E)),
    ("Security", "CSRF · Rate Limiting · Password Hashing · Role-Based Access · Audit Logging · Environment Secrets", DARK_BLUE),
    ("Additional Skills", "OCR (Tesseract) · SMTP Email · Web Scraping · Data Migration · Excel Import/Export · REST APIs", ACCENT),
    ("Languages", "Python (expert) · SQL (advanced) · JavaScript (intermediate) · HTML/CSS (expert) · Bash (intermediate)", PRIMARY),
]
for i, (title, desc, color) in enumerate(skills):
    row, col = divmod(i, 2)
    x = Inches(0.5) + col * Inches(6.3)
    y = Inches(1.2) + row * Inches(1.35)
    c = add_rnd(slide, x, y, Inches(5.9), Inches(1.15), fill=CARD_BG)
    add_shape(slide, x, y, Inches(0.05), Inches(1.15), fill=color)
    tb(slide, x + Inches(0.2), y + Inches(0.08), Inches(5.4), Inches(0.3), title, 13, color, True)
    tb(slide, x + Inches(0.2), y + Inches(0.42), Inches(5.4), Inches(0.6), desc, 10, GRAY)

# ════════════════════════════════════════════════════════════
# SLIDE 12 — Stats Summary
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, W, H, fill=LIGHT)
section_header(slide, "Project Stats at a Glance", "Built over 2+ years for a real business in Abu Dhabi")

big_stats = [
    ("19,300+", "Python Lines", "Single routes.py + 2,100+ in blueprint routes + 3,900 in pdf_service.py + services", PRIMARY),
    ("1,200+", "Git Commits", "Clean linear history, semantic commit messages, feature branches", ACCENT),
    ("50+", "Database Tables", "12 modules, fully normalized, with foreign key relationships", GOLD),
    ("20+", "PDF Types", "Professional ReportLab documents — invoices, SOA, slips, vouchers, cheques", RGBColor(0x7C, 0x3A, 0xED)),
    ("12", "Integrated Modules", "HR, Fleet, Customer, Supplier, Accounts, Documents, Portals, Reports", RGBColor(0xE5, 0x3E, 0x3E)),
    ("5+", "User Portals", "Admin · Accounts · Supplier · Driver · Field Staff — each with separate login", DARK_BLUE),
]
for i, (num, title, desc, color) in enumerate(big_stats):
    row, col = divmod(i, 3)
    x = Inches(0.5) + col * Inches(4.2)
    y = Inches(1.3) + row * Inches(2.7)
    c = add_rnd(slide, x, y, Inches(3.9), Inches(2.3), fill=WHITE)
    add_shape(slide, x, y, Inches(3.9), Inches(0.06), fill=color)
    tb(slide, x + Inches(0.15), y + Inches(0.2), Inches(3.6), Inches(0.65), num, 36, color, True, PP_ALIGN.CENTER)
    tb(slide, x + Inches(0.15), y + Inches(0.9), Inches(3.6), Inches(0.35), title, 14, DARK, True, PP_ALIGN.CENTER)
    tb(slide, x + Inches(0.15), y + Inches(1.3), Inches(3.6), Inches(0.85), desc, 10, GRAY, False, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 13 — Contact / CTA
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, W, H, fill=DARK)
add_shape(slide, 0, 0, Inches(0.12), H, fill=PRIMARY)
# Decorative
s1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(5), Inches(4), Inches(4))
s1.fill.solid(); s1.fill.fore_color.rgb = NAVY; s1.line.fill.background()

tb(slide, Inches(1), Inches(1.8), Inches(10), Inches(0.4), "Let's Build Something Great Together", 40, WHITE, True)
add_shape(slide, Inches(1), Inches(2.4), Inches(2.5), Inches(0.03), fill=GOLD)
tb(slide, Inches(1), Inches(2.8), Inches(10), Inches(0.5), "Full-Stack Python Developer | ERP Specialist | Ready for New Challenges", 16, GOLD)

# Contact box
cb = add_rnd(slide, Inches(1), Inches(3.8), Inches(7), Inches(2.8), fill=NAVY)
tb(slide, Inches(1.3), Inches(4.0), Inches(6), Inches(0.3), "Contact Information", 14, WHITE, True)
contacts = [
    ("📞", "+971 50-122-4963"),
    ("📧", "info@currentlinktgc.com"),
    ("🌐", "www.currentlinkgc.com"),
    ("📍", "Mussaffah, Abu Dhabi, UAE"),
]
for i, (icon, info) in enumerate(contacts):
    y = Inches(4.4) + i * Inches(0.5)
    tb(slide, Inches(1.3), y, Inches(0.4), Inches(0.35), icon, 14, GOLD)
    tb(slide, Inches(1.8), y, Inches(5), Inches(0.35), info, 13, WHITE)

# Right side tech mentions
right_box = add_rnd(slide, Inches(8.5), Inches(3.8), Inches(4.2), Inches(2.8), fill=CARD_BG)
tb(slide, Inches(8.8), Inches(4.0), Inches(3.6), Inches(0.3), "Core Competencies", 13, GOLD, True)
competencies = [
    "✓ Full-Stack Python Development",
    "✓ ERP System Architecture",
    "✓ Database Design (SQL)",
    "✓ PDF Generation Engineering",
    "✓ Linux Server Administration",
    "✓ Git & DevOps Practices",
    "✓ Problem Solving & Debugging",
    "✓ End-to-End Product Ownership",
]
bullets(slide, Inches(8.8), Inches(4.4), Inches(3.6), Inches(2.0), competencies, 11, GRAY, 5)

tb(slide, Inches(1), Inches(6.9), Inches(10), Inches(0.3), "Available for remote · Freelance · Full-Time | Python Flask · PostgreSQL · Full-Stack ERP Development", 11, GRAY)

# ── Save ──
output_path = os.path.join(os.environ['USERPROFILE'], "Desktop", "Current_Link_ERP_Pitch_Deck.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Size: {os.path.getsize(output_path)/1024:.0f} KB | Slides: {len(prs.slides)}")
